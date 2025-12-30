"""
Content extraction services
"""
from typing import Dict, Any
import io
from loguru import logger


async def extract_content_from_file(file_content: bytes, filename: str) -> Dict[str, Any]:
    """
    Extract content from uploaded file
    Supports: PDF, DOCX, PPTX, TXT
    """
    # Determine file type
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    try:
        if ext == 'pdf':
            return await _extract_from_pdf(file_content, filename)
        elif ext == 'docx':
            return await _extract_from_docx(file_content, filename)
        elif ext == 'pptx':
            return await _extract_from_pptx(file_content, filename)
        elif ext == 'txt':
            return await _extract_from_txt(file_content, filename)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    
    except Exception as e:
        logger.error(f"Error extracting content from {filename}: {e}")
        raise ValueError(f"Failed to extract content from {filename}: {str(e)}")


async def _extract_from_pdf(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from PDF using pdfplumber"""
    import pdfplumber
    
    text_parts = []
    metadata = {}
    
    try:
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            metadata['page_count'] = len(pdf.pages)
            
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            
            # Try to extract title from metadata
            if pdf.metadata:
                title = pdf.metadata.get('Title') or pdf.metadata.get('title')
                if title:
                    metadata['original_title'] = title
        
        full_text = '\n\n'.join(text_parts)
        
        if not full_text.strip():
            raise ValueError("PDF appears to be empty or contains no extractable text")
        
        metadata.update({
            'word_count': len(full_text.split()),
            'character_count': len(full_text),
            'file_size': len(file_content)
        })
        
        return {
            "text": full_text,
            "title": metadata.get('original_title') or filename.rsplit('.', 1)[0],
            "source_type": "pdf",
            "metadata": metadata
        }
    
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        raise ValueError(f"Failed to extract PDF content: {str(e)}")


async def _extract_from_docx(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from DOCX using python-docx"""
    from docx import Document
    
    try:
        doc = Document(io.BytesIO(file_content))
        
        # Extract paragraphs
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
        full_text = '\n\n'.join(text_parts)
        
        if not full_text.strip():
            raise ValueError("DOCX appears to be empty")
        
        # Extract core properties if available
        metadata = {
            'paragraph_count': len(text_parts),
            'word_count': len(full_text.split()),
            'character_count': len(full_text),
            'file_size': len(file_content)
        }
        
        # Try to get title from document properties
        title = filename.rsplit('.', 1)[0]
        if hasattr(doc.core_properties, 'title') and doc.core_properties.title:
            title = doc.core_properties.title
            metadata['original_title'] = title
        
        return {
            "text": full_text,
            "title": title,
            "source_type": "docx",
            "metadata": metadata
        }
    
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        raise ValueError(f"Failed to extract DOCX content: {str(e)}")


async def _extract_from_pptx(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from PPTX using python-pptx"""
    from pptx import Presentation
    
    try:
        prs = Presentation(io.BytesIO(file_content))
        
        text_parts = []
        slide_count = 0
        
        for slide in prs.slides:
            slide_count += 1
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"--- Slide {slide_count} ---\n" + '\n'.join(slide_text))
        
        full_text = '\n\n'.join(text_parts)
        
        if not full_text.strip():
            raise ValueError("PPTX appears to be empty or contains no extractable text")
        
        metadata = {
            'slide_count': slide_count,
            'word_count': len(full_text.split()),
            'character_count': len(full_text),
            'file_size': len(file_content)
        }
        
        return {
            "text": full_text,
            "title": filename.rsplit('.', 1)[0],
            "source_type": "pptx",
            "metadata": metadata
        }
    
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        raise ValueError(f"Failed to extract PPTX content: {str(e)}")


async def _extract_from_txt(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from TXT file"""
    try:
        # Try UTF-8 first, then fallback to latin-1
        try:
            text = file_content.decode('utf-8')
        except UnicodeDecodeError:
            text = file_content.decode('latin-1')
        
        if not text.strip():
            raise ValueError("TXT file appears to be empty")
        
        metadata = {
            'word_count': len(text.split()),
            'character_count': len(text),
            'file_size': len(file_content),
            'line_count': len(text.splitlines())
        }
        
        return {
            "text": text,
            "title": filename.rsplit('.', 1)[0],
            "source_type": "txt",
            "metadata": metadata
        }
    
    except Exception as e:
        logger.error(f"TXT extraction error: {e}")
        raise ValueError(f"Failed to extract TXT content: {str(e)}")


async def extract_content_from_url(url: str) -> Dict[str, Any]:
    """
    Extract content from URL using web scraping
    """
    import aiohttp
    from bs4 import BeautifulSoup
    
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=30) as response:
                if response.status != 200:
                    raise ValueError(f"Failed to fetch URL: HTTP {response.status}")
                
                html = await response.text()
        
        # Parse HTML
        soup = BeautifulSoup(html, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()
        
        # Extract title
        title = "Untitled"
        if soup.title and soup.title.string:
            title = soup.title.string.strip()
        
        # Extract main content
        # Try to find main content area
        main_content = soup.find('main') or soup.find('article') or soup.find('div', class_='content')
        
        if main_content:
            text = main_content.get_text(separator='\n', strip=True)
        else:
            text = soup.get_text(separator='\n', strip=True)
        
        # Clean up excessive newlines
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        text = '\n\n'.join(lines)
        
        if not text.strip():
            raise ValueError("No extractable text found on the webpage")
        
        metadata = {
            "url": url,
            "word_count": len(text.split()),
            "character_count": len(text)
        }
        
        return {
            "text": text,
            "title": title,
            "source_type": "url",
            "metadata": metadata
        }
    
    except Exception as e:
        logger.error(f"URL extraction error for {url}: {e}")
        raise ValueError(f"Failed to extract content from URL: {str(e)}")
